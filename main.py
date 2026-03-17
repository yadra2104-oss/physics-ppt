from pptx import Presentation
from pptx.util import Pt

def create_physics_presentation():
    prs = Presentation()

    # Функція для додавання слайда
    def add_slide(title_text, content_lines):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.size = Pt(40)
        title.text_frame.paragraphs[0].font.bold = True
        
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        for line in content_lines:
            p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(24)

    # --- Слайди ---
    # Слайд 1 (Титульний)
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Провідники та діелектрики в електричному полі"
    subtitle.text = "Основи електростатики\n\nВиконав: Учень 10-А класу"

    # Слайд 2
    add_slide(
        "Електричні властивості речовин",
        [
            "Всі речовини поділяються на:",
            "1. Провідники (метали, електроліти).",
            "2. Діелектрики (ізолятори).",
            "3. Напівпровідники.",
            "Основна відмінність — наявність вільних носіїв заряду."
        ]
    )

    # Слайд 3
    add_slide(
        "Провідники в електричному полі",
        [
            "Провідники — речовини з вільними зарядженими частинками.",
            "Приклади: метали (електрони), розчини солей (іони).",
            "У зовнішньому полі відбувається електростатична індукція.",
            "Напруженість поля всередині провідника E=0."
        ]
    )

    # Слайд 4
    add_slide(
        "Діелектрики в електричному полі",
        [
            "Діелектрики — відсутні вільні носії заряду.",
            "Приклади: скло, гума, порцеляна, повітря.",
            "Заряди в діелектрику є зв'язаними.",
            "Під дією поля відбувається поляризація."
        ]
    )

    # Слайд 5
    add_slide(
        "Висновки",
        [
            "1. Провідники мають вільні заряди.",
            "2. Діелектрики мають зв'язані заряди.",
            "3. Поле всередині провідника дорівнює нулю.",
            "4. У діелектрику поле послаблюється."
        ]
    )

    # Збереження
    prs.save('presentation.pptx')
    print("Файл presentation.pptx успішно створено!")

# Запуск
if __name__ == "__main__":
    create_physics_presentation()
